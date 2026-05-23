from dataclasses import dataclass
from pathlib import Path
from typing import List

from fastapi import HTTPException

from app.services.pdf.pdf_loader import extract_pdf_text


@dataclass
class PageText:
    page: int
    text: str
    used_ocr: bool = False


def extract_document_text(path: Path) -> List[PageText]:
    extension = path.suffix.lower()
    if extension == ".pdf":
        return extract_pdf_text(path)
    if extension == ".docx":
        return extract_docx_text(path)
    if extension == ".pptx":
        return extract_pptx_text(path)
    raise HTTPException(status_code=400, detail=f"Unsupported document type: {extension or 'unknown'}")


def extract_docx_text(path: Path) -> List[PageText]:
    try:
        from docx import Document
    except ImportError as exc:
        raise HTTPException(status_code=500, detail="Word document support is not installed.") from exc

    document = Document(path)
    parts: List[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    text = "\n".join(parts).strip()
    return [PageText(page=1, text=text)] if text else []


def extract_pptx_text(path: Path) -> List[PageText]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise HTTPException(status_code=500, detail="PowerPoint document support is not installed.") from exc

    presentation = Presentation(path)
    pages: List[PageText] = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))

        text = "\n".join(parts).strip()
        pages.append(PageText(page=index, text=text))

    return pages
